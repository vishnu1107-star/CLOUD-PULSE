import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    BG_WHITE = RGBColor(255, 255, 255)
    CARD_BG = RGBColor(248, 250, 252)
    CARD_BORDER = RGBColor(226, 232, 240)
    
    TEXT_DARK = RGBColor(15, 23, 42)
    TEXT_MUTED = RGBColor(71, 85, 105)
    
    PRIMARY_BLUE = RGBColor(29, 78, 216)
    ACCENT_TEAL = RGBColor(13, 148, 136)
    ACCENT_INDIGO = RGBColor(79, 70, 229)
    
    blank_layout = prs.slide_layouts[6]
    
    def set_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_WHITE
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text="CLOUDPULSE INNOVATION PROJECT"):
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_TEAL
        
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_DARK

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_background(slide1)
    
    card1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.0), Inches(11.333), Inches(5.5))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = PRIMARY_BLUE
    card1.line.width = Pt(2)
    
    accent_bar = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.0), Inches(11.333), Inches(0.15))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = PRIMARY_BLUE
    accent_bar.line.fill.background()

    tbox = slide1.shapes.add_textbox(Inches(1.5), Inches(1.6), Inches(10.333), Inches(1.8))
    tf = tbox.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = "⚡ CloudPulse"
    p0.font.size = Pt(46)
    p0.font.bold = True
    p0.font.color.rgb = PRIMARY_BLUE
    
    p1 = tf.add_paragraph()
    p1.text = "Autonomous Multi-Cloud Idle Reclamation & Instant Hydration Engine"
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_DARK
    p1.space_before = Pt(6)

    desc_box = slide1.shapes.add_textbox(Inches(1.5), Inches(3.6), Inches(10.333), Inches(2.2))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    
    p3 = tf_desc.paragraphs[0]
    p3.text = "A novel zero-downtime FinOps architecture combining multi-signal idle telemetry, ghost resource reaper, 3-second developer re-activation, and real-time carbon offset accounting."
    p3.font.size = Pt(15)
    p3.font.color.rgb = TEXT_MUTED
    p3.space_after = Pt(16)
    
    p4 = tf_desc.add_paragraph()
    p4.text = "PROPOSED INNOVATION • DETAILED PLAN • TECH STACK"
    p4.font.size = Pt(13)
    p4.font.bold = True
    p4.font.color.rgb = ACCENT_TEAL

    # -------------------------------------------------------------
    # SLIDE 2: Problem Statement
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_background(slide2)
    add_header(slide2, "Problem Statement: The $17B Cloud Idle Crisis")
    
    problems = [
        {
            "num": "01",
            "title": "Severe Idle Dev Burn",
            "subtitle": "40%+ Budget Wasted Off-Hours",
            "points": [
                "Non-production VMs (AWS EC2, GCP Compute) and Kubernetes clusters run 24/7 unnecessarily.",
                "Over 68% of total weekly hours are completely idle with zero developer activity."
            ]
        },
        {
            "num": "02",
            "title": "Silent Ghost Assets",
            "subtitle": "Orphaned Storage & IP Drain",
            "points": [
                "Unattached EBS/GCP disks, orphaned Elastic IPs, and idle NAT gateways silently accumulate monthly charges.",
                "Cloud providers charge continuously regardless of attached instance state."
            ]
        },
        {
            "num": "03",
            "title": "Flawed Advisory FinOps",
            "subtitle": "Alert Fatigue & Manual Friction",
            "points": [
                "Legacy tools only output passive PDF/email recommendations without executing safe actions.",
                "Manual cleanup causes developer friction and fear of accidentally breaking dependencies."
            ]
        }
    ]
    
    left_positions = [Inches(0.8), Inches(4.8), Inches(8.8)]
    card_width = Inches(3.733)
    card_height = Inches(4.8)
    
    for i, prob in enumerate(problems):
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_positions[i], Inches(1.8), card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1.5)
        
        tab = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_positions[i], Inches(1.8), card_width, Inches(0.12))
        tab.fill.solid()
        tab.fill.fore_color.rgb = PRIMARY_BLUE if i == 0 else (ACCENT_TEAL if i == 1 else ACCENT_INDIGO)
        tab.line.fill.background()

        tb = slide2.shapes.add_textbox(left_positions[i] + Inches(0.25), Inches(2.1), card_width - Inches(0.5), card_height - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = prob["num"]
        p0.font.size = Pt(26)
        p0.font.bold = True
        p0.font.color.rgb = PRIMARY_BLUE
        
        p1 = tf.add_paragraph()
        p1.text = prob["title"]
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_DARK
        
        p_sub = tf.add_paragraph()
        p_sub.text = prob["subtitle"]
        p_sub.font.size = Pt(12)
        p_sub.font.bold = True
        p_sub.font.color.rgb = ACCENT_TEAL
        p_sub.space_after = Pt(12)
        
        for pt in prob["points"]:
            p_pt = tf.add_paragraph()
            p_pt.text = "• " + pt
            p_pt.font.size = Pt(13)
            p_pt.font.color.rgb = TEXT_MUTED
            p_pt.space_after = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 3: Proposed Solution
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_background(slide3)
    add_header(slide3, "Proposed Solution: The Novel CloudPulse Engine Architecture")
    
    solutions = [
        ("1. Multi-Signal Telemetry Fusion", "Combines CPU (<2%), Network (<10KB/s), IOPS, and active DB/HTTP socket checks over rolling 30-min windows to prevent false positives."),
        ("2. Safe Tag-Aware Governance", "Dynamically exempts critical production workloads (<code>Environment: Production</code>) while governing staging/dev resources."),
        ("3. Instant-Warm Hydration Protocol", "3-second developer re-activation via Next.js Web Portal or Slack Slash Command (<code>/cloudpulse wakeup staging</code>) with zero state loss."),
        ("4. Autonomous Ghost Resource Reaper", "Purges unattached EBS volumes, orphan EIPs, and idle load balancers with automated 30-day rollback snapshots."),
        ("5. Auditable FinOps & Carbon Ledger", "Real-time interactive dashboard tracking dollar savings ($) and verified carbon footprint reduction (kg CO₂ emissions avoided).")
    ]
    
    sol_y_start = Inches(1.8)
    for i, (title, desc) in enumerate(solutions):
        y_pos = sol_y_start + Inches(i * 0.98)
        
        strip = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_pos, Inches(11.733), Inches(0.85))
        strip.fill.solid()
        strip.fill.fore_color.rgb = CARD_BG
        strip.line.color.rgb = PRIMARY_BLUE if i % 2 == 0 else ACCENT_TEAL
        strip.line.width = Pt(1)
        
        tb = slide3.shapes.add_textbox(Inches(1.0), y_pos + Inches(0.08), Inches(11.333), Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title + "  —  "
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_BLUE
        
        run = p.add_run()
        run.text = desc
        run.font.size = Pt(13)
        run.font.bold = False
        run.font.color.rgb = TEXT_DARK

    # -------------------------------------------------------------
    # SLIDE 4: Tech Stack
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_background(slide4)
    add_header(slide4, "Technology Stack & Layering")
    
    stacks = [
        {
            "category": "Frontend & Web3 Portal",
            "techs": "Next.js 14, React 18, TailwindCSS, Web3.js / Solana Py, Recharts Analytics",
            "desc": "Executive web dashboard for telemetry, ghost resource sweeping, 1-click hydration, and Web3 carbon credit minting."
        },
        {
            "category": "AI Telemetry & Backend Engine",
            "techs": "FastAPI (Python 3.11), Time-Series AI Engine (Scikit/Prophet), Async SQLAlchemy, PostgreSQL",
            "desc": "High-concurrency REST engine handling predictive AI telemetry evaluation, scheduled job dispatch, and secure webhook verification."
        },
        {
            "category": "Multi-Cloud & VEGA Hardware",
            "techs": "C-DAC VEGA RISC-V SoC Board, AWS Boto3 SDK, GCP Compute API, K8s SDK, Docker",
            "desc": "Physical VEGA RISC-V Edge Hardware Gateway paired with native cloud drivers executing zero-downtime pause/scale commands."
        },
        {
            "category": "DevOps, ChatOps & Web3 Ledger",
            "techs": "Slack Webhooks API, On-Chain Carbon Ledger (Solana/Polygon), Docker Compose, OpenAPI Specs",
            "desc": "Instant ChatOps re-activation (`/cloudpulse wakeup`), on-chain proof-of-green compute ledger, and containerized deployment."
        }
    ]
    
    box_positions = [
        (Inches(0.8), Inches(1.8)),
        (Inches(6.8), Inches(1.8)),
        (Inches(0.8), Inches(4.5)),
        (Inches(6.8), Inches(4.5))
    ]
    
    for i, stack in enumerate(stacks):
        x, y = box_positions[i]
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.733), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1.5)
        
        stripe = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(0.12), Inches(2.4))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = ACCENT_TEAL if i % 2 == 0 else PRIMARY_BLUE
        stripe.line.fill.background()

        tb = slide4.shapes.add_textbox(x + Inches(0.25), y + Inches(0.15), Inches(5.3), Inches(2.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = stack["category"]
        p0.font.size = Pt(17)
        p0.font.bold = True
        p0.font.color.rgb = PRIMARY_BLUE
        
        p1 = tf.add_paragraph()
        p1.text = "Core Tech: " + stack["techs"]
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = ACCENT_TEAL
        p1.space_after = Pt(6)
        
        p2 = tf.add_paragraph()
        p2.text = stack["desc"]
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 5: Implementation Plan & Strategic Value
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_background(slide5)
    add_header(slide5, "Clear Implementation Plan & Expected Outcomes")
    
    phases = [
        ("Phase 1: Discovery & Telemetry Sync (W1-W2)", "Deploy multi-cloud provider drivers, configure tag-aware asset inventory, and baseline idle metrics."),
        ("Phase 2: Sweeping & ChatOps Hydration (W3-W4)", "Enable dry-run ghost resource reaper, integrate K8s auto-scaler to 0, and launch Slack Slash commands."),
        ("Phase 3: Enforcement & Carbon Ledger (W5-W6)", "Activate automated lifecycle policies, real-time cost reclamation analytics, and CO2 offset ledger.")
    ]
    
    for i, (p_title, p_desc) in enumerate(phases):
        y_pos = Inches(1.8) + Inches(i * 1.0)
        
        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_pos, Inches(11.733), Inches(0.85))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = PRIMARY_BLUE
        card.line.width = Pt(1.5)
        
        tb = slide5.shapes.add_textbox(Inches(1.0), y_pos + Inches(0.08), Inches(11.333), Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = p_title + "  —  "
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_BLUE
        
        run = p.add_run()
        run.text = p_desc
        run.font.size = Pt(12)
        run.font.color.rgb = TEXT_DARK

    summary_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.9), Inches(11.733), Inches(1.4))
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = RGBColor(239, 246, 255)
    summary_box.line.color.rgb = ACCENT_TEAL
    summary_box.line.width = Pt(1.5)
    
    tb_sum = slide5.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(11.333), Inches(1.2))
    tf_sum = tb_sum.text_frame
    tf_sum.word_wrap = True
    
    ps0 = tf_sum.paragraphs[0]
    ps0.text = "PROJECTED STRATEGIC IMPACT & ROI"
    ps0.font.size = Pt(13)
    ps0.font.bold = True
    ps0.font.color.rgb = ACCENT_TEAL
    ps0.space_after = Pt(4)
    
    ps1 = tf_sum.add_paragraph()
    ps1.text = "• 55% - 65% Reduction in non-production cloud spend within 30 days of deployment."
    ps1.font.size = Pt(12)
    ps1.font.bold = True
    ps1.font.color.rgb = TEXT_DARK
    
    ps2 = tf_sum.add_paragraph()
    ps2.text = "• 100% Reclamation of ghost EBS/GCP volumes and orphan IPs with 0% developer workflow disruption."
    ps2.font.size = Pt(12)
    ps2.font.bold = True
    ps2.font.color.rgb = TEXT_DARK

    # Output to both Root Workspace and cloudpulse directory
    paths = [
        os.path.join(r"C:\Users\dELL\OneDrive\Desktop\main-2", "CloudPulse_Presentation_Updated.pptx"),
        os.path.join(r"C:\Users\dELL\OneDrive\Desktop\main-2\cloudpulse", "CloudPulse_Presentation_Updated.pptx"),
        os.path.join(r"C:\Users\dELL\OneDrive\Desktop\main-2", "CloudPulse_Presentation_Final.pptx"),
        os.path.join(r"C:\Users\dELL\OneDrive\Desktop\main-2\cloudpulse", "CloudPulse_Presentation_Final.pptx")
    ]
    for p in paths:
        try:
            prs.save(p)
            print(f"Saved PPTX to {p}")
        except Exception as e:
            print(f"Skipped locked file {p}: {e}")

if __name__ == "__main__":
    create_presentation()
